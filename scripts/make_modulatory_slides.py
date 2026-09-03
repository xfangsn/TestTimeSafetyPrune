"""Editable slide deck: Driving vs Modulatory weights, and BLADE targets the
modulatory ones. All native pptx shapes/text -> fully editable in PowerPoint /
Google Slides / Keynote. Output: figures/blade_modulatory_slides.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIG = Path("figures"); FIG.mkdir(exist_ok=True)

# palette
INK = RGBColor(0x22, 0x22, 0x22)
BLUE = RGBColor(0x1D, 0x4E, 0x89); BLUE_F = RGBColor(0xD5, 0xE5, 0xF7)
AMB = RGBColor(0xC8, 0x7A, 0x00); AMB_F = RGBColor(0xFB, 0xE7, 0xC6)
GRN = RGBColor(0x2A, 0x7F, 0x4F); GRN_F = RGBColor(0xD9, 0xEF, 0xE1)
GREY = RGBColor(0x8A, 0x8A, 0x8A); WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def box(sl, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill if fill else WHITE
    if fill is None:
        sp.fill.background()
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def txt(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for i, (t, sz, b, c) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c
        r.font.name = "Arial"
    return tb


def shape_text(sp, runs, align=PP_ALIGN.CENTER):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (t, sz, b, c) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = "Arial"


# ===================== SLIDE 1: concept diagram =====================
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, fill=WHITE)
txt(s, 0.5, 0.28, 12.3, 0.7,
    [("Two kinds of weights in an LLM — BLADE edits the modulatory ones", 26, True, INK)])
txt(s, 0.5, 0.98, 12.3, 0.5,
    [("Control a behavior by zeroing a few gain/gate weights, without touching what the model can do",
      15, False, GREY)])

# --- DRIVING lane (blue) ---
box(s, 0.5, 1.8, 7.4, 2.05, fill=BLUE_F, line=BLUE, lw=1.5)
txt(s, 0.7, 1.9, 5, 0.4, [("DRIVING weights  (W_drive)", 15, True, BLUE)])
# input dots -> content node -> output, with thick arrows
for i, yy in enumerate([2.45, 2.9, 3.35]):
    c = box(s, 0.9, yy, 0.32, 0.32, fill=WHITE, line=BLUE, lw=1.2, shape=MSO_SHAPE.OVAL)
a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.4), Inches(2.55), Inches(1.7), Inches(0.75))
a.fill.solid(); a.fill.fore_color.rgb = BLUE; a.line.fill.background(); a.shadow.inherit = False
cn = box(s, 3.25, 2.35, 1.9, 1.2, fill=WHITE, line=BLUE, lw=1.4)
shape_text(cn, [("signal /\ncontent", 12, True, BLUE)])
a2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.3), Inches(2.6), Inches(1.2), Inches(0.65))
a2.fill.solid(); a2.fill.fore_color.rgb = BLUE; a2.line.fill.background(); a2.shadow.inherit = False
out = box(s, 6.6, 2.4, 1.1, 1.1, fill=WHITE, line=BLUE, lw=1.4, shape=MSO_SHAPE.OVAL)
shape_text(out, [("out", 12, True, BLUE)])
txt(s, 0.7, 3.42, 7, 0.35,
    [("carry the signal  →  WHAT the model can do (capability)", 12.5, False, INK)])

# --- MODULATORY gate (amber) sitting on the output path ---
gate = box(s, 5.55, 2.02, 0.7, 1.9, fill=AMB_F, line=AMB, lw=1.6, shape=MSO_SHAPE.RECTANGLE)
knob = box(s, 5.62, 2.55, 0.56, 0.56, fill=AMB, line=AMB, lw=1, shape=MSO_SHAPE.OVAL)
shape_text(knob, [("⦿", 14, True, WHITE)])
box(s, 8.15, 1.8, 4.65, 2.05, fill=AMB_F, line=AMB, lw=1.5)
txt(s, 8.35, 1.9, 4.3, 0.4, [("MODULATORY weights  (W_mod)", 15, True, AMB)])
txt(s, 8.35, 2.35, 4.3, 1.4,
    [("a gain / gate on the path", 13, True, AMB),
     ("→ WHETHER a behavior is expressed (how it behaves)", 12.5, False, INK),
     ("sparse — a few residual-writer edges", 12, False, GREY)])

# --- BLADE callout (green) pointing at the knob ---
bl = box(s, 3.9, 4.35, 5.55, 1.0, fill=GRN_F, line=GRN, lw=1.6)
shape_text(bl, [("BLADE  (gradient-free, forward-only)", 14, True, GRN),
                ("scores  s_ij = [ r · W · Δμ ]₊  and zeros ~0.2–1% modulatory writer weights",
                 12, False, INK)])
ar = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(5.7), Inches(3.95), Inches(0.4), Inches(0.42))
ar.fill.solid(); ar.fill.fore_color.rgb = GRN; ar.line.fill.background(); ar.shadow.inherit = False

# --- bottom contrast table ---
yb = 5.6
d = box(s, 0.5, yb, 6.05, 1.6, fill=WHITE, line=BLUE, lw=1.4)
txt(s, 0.7, yb + 0.08, 5.7, 1.5,
    [("Zero DRIVING weights", 13, True, BLUE),
     ("• dense, distributed   • carry content", 11.5, False, INK),
     ("• → capability collapses (ppl ↑↑)", 11.5, False, INK),
     ("• neuro analogy: driving synapse (sets firing)", 11.5, False, GREY)])
m = box(s, 6.78, yb, 6.05, 1.6, fill=WHITE, line=AMB, lw=1.4)
txt(s, 6.98, yb + 0.08, 5.7, 1.5,
    [("Zero MODULATORY weights  ← BLADE", 13, True, AMB),
     ("• sparse, few edges   • gain / gating of a behavior", 11.5, False, INK),
     ("• → behavior OFF, capability intact (ppl ≈ 0)", 11.5, False, INK),
     ("• neuro analogy: modulatory synapse (tunes gain)", 11.5, False, GREY)])

# ===================== SLIDE 2: evidence =====================
s2 = prs.slides.add_slide(BLANK)
box(s2, 0, 0, 13.333, 7.5, fill=WHITE)
txt(s2, 0.5, 0.3, 12.3, 0.7,
    [("Evidence: the weights BLADE finds are modulatory, not driving", 25, True, INK)])
txt(s2, 0.5, 1.0, 12.3, 0.45,
    [("Zeroing them removes the behavior's expression while capability is preserved", 15, False, GREY)])

rows = [
    ("Refusal (Llama-3.2-3B, OOD HarmBench)",
     "zero ~0.07% writer weights in ONE mid-layer  →  refusal 0.91 → 0.09;  WikiText Δppl +0.3%;  6-task downstream unchanged.  The model still CAN produce the content — the gate is gone, the capability is not."),
    ("Sycophancy (A/B preference)",
     "0.88 → 0.50 (no preference) at Δppl ≈ +1%,  downstream +0.1pp.  Behavior expression removed; general ability intact."),
    ("Contrast: Weight-Steering (adds a task vector)",
     "matches BLADE on OOD but costs −3.8pp downstream — a driving-scale edit that moves capability, not just the gate."),
    ("Synapse verdict (independent review: codex + kimi)",
     "the selected weights act as gain / gating (modulatory). Removing them changes behavior, not the underlying computation."),
]
y = 1.7
for head, body in rows:
    box(s2, 0.5, y, 12.33, 1.18, fill=GRN_F if head.startswith("Synapse") else WHITE,
        line=GRN if head.startswith("Synapse") else GREY, lw=1.2)
    txt(s2, 0.75, y + 0.1, 11.9, 1.0,
        [(head, 14, True, GRN if head.startswith("Synapse") else INK),
         (body, 12.5, False, INK)])
    y += 1.32

out = FIG / "blade_modulatory_slides.pptx"
prs.save(str(out))
print(f"saved {out}")
